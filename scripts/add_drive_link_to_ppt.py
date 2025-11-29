"""
Add Google Drive link to first and last slides of PowerPoint
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

# Load presentation
prs = Presentation(r'C:\Users\USER\Documents\10Alatytics-2025\presentation\10Alytics_Governance_Growth_Gap.pptx')

print(f"Total slides: {len(prs.slides)}\n")

# Add submission links to slide 1 (after title)
slide_1 = prs.slides[0]

# Remove existing textboxes at bottom if any (from previous run)
shapes_to_remove = []
for shape in slide_1.shapes:
    if hasattr(shape, "top") and shape.top > Inches(6):
        shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    sp = shape.element
    sp.getparent().remove(sp)

# Add new text box at bottom
left = Inches(0.5)
top = Inches(6.5)
width = Inches(9)
height = Inches(1)

textbox = slide_1.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "📦 Full Submission Package (ZIP - 13MB - Please Extract)"
p.font.size = Pt(11)
p.font.bold = True
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

p2 = text_frame.add_paragraph()
p2.text = "https://drive.google.com/file/d/1lhvmFr-v9KLngbqpYUu3as_HXbGbfULp/view?usp=drivesdk"
p2.font.size = Pt(9)
p2.font.name = 'Calibri'
p2.alignment = PP_ALIGN.CENTER

p3 = text_frame.add_paragraph()
p3.text = "GitHub: https://github.com/Basit-Balogun10/africa-governance-growth-gap-analysis"
p3.font.size = Pt(8)
p3.font.name = 'Calibri'
p3.alignment = PP_ALIGN.CENTER

print(f"✓ Added submission links to slide 1")

# Add to last slide as well
last_slide = prs.slides[-1]

# Remove existing textboxes at bottom if any
shapes_to_remove_last = []
for shape in last_slide.shapes:
    if hasattr(shape, "top") and shape.top > Inches(5.5):
        shapes_to_remove_last.append(shape)

for shape in shapes_to_remove_last:
    sp = shape.element
    sp.getparent().remove(sp)

textbox_last = last_slide.shapes.add_textbox(left, Inches(6), width, Inches(1.2))
text_frame_last = textbox_last.text_frame
text_frame_last.word_wrap = True

p_last = text_frame_last.paragraphs[0]
p_last.text = "📦 Full Working Files Submission (ZIP - 13MB - Extract Required)"
p_last.font.size = Pt(11)
p_last.font.bold = True
p_last.font.name = 'Calibri'
p_last.alignment = PP_ALIGN.CENTER

p_last2 = text_frame_last.add_paragraph()
p_last2.text = "Google Drive: https://drive.google.com/file/d/1lhvmFr-v9KLngbqpYUu3as_HXbGbfULp/view?usp=drivesdk"
p_last2.font.size = Pt(9)
p_last2.font.name = 'Calibri'
p_last2.alignment = PP_ALIGN.CENTER

p_last3 = text_frame_last.add_paragraph()
p_last3.text = "GitHub Repo: https://github.com/Basit-Balogun10/africa-governance-growth-gap-analysis"
p_last3.font.size = Pt(8)
p_last3.font.name = 'Calibri'
p_last3.alignment = PP_ALIGN.CENTER

p_last4 = text_frame_last.add_paragraph()
p_last4.text = "Email: basitbalogun10@gmail.com"
p_last4.font.size = Pt(9)
p_last4.font.name = 'Calibri'
p_last4.alignment = PP_ALIGN.CENTER

print(f"✓ Added submission links to slide {len(prs.slides)}")

# Save
prs.save(r'C:\Users\USER\Documents\10Alatytics-2025\presentation\10Alytics_Governance_Growth_Gap.pptx')
print(f"\n✅ PowerPoint updated with Google Drive links!")
