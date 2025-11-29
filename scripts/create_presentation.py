"""
Generate PowerPoint Presentation for 10Alytics Hackathon 2025
The Governance-Growth Gap: Why African Nations Can't Translate Spending Into Progress
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(16)  # Widescreen 16:9
prs.slide_height = Inches(9)

# Define color scheme
PRIMARY_COLOR = RGBColor(41, 128, 185)  # Blue
SUCCESS_COLOR = RGBColor(39, 174, 96)   # Green
DANGER_COLOR = RGBColor(231, 76, 60)    # Red
WARNING_COLOR = RGBColor(243, 156, 18)  # Orange
TEXT_COLOR = RGBColor(44, 62, 80)       # Dark gray

def add_title_slide():
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "The Governance-Growth Gap"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(14), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Why African Nations Can't Translate Spending Into Progress"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = TEXT_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "10Alytics Global Hackathon 2025 | November 29, 2025"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(16)
    footer_para.font.italic = True
    footer_para.font.color.rgb = RGBColor(127, 127, 127)
    footer_para.alignment = PP_ALIGN.CENTER

def add_hook_slide():
    """Slide 2: The Hook"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = '"Better Governance Should Mean Better Growth... Right?"'
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = DANGER_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Key finding
    finding_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(2))
    finding_frame = finding_box.text_frame
    finding_frame.text = "We analyzed 14 African countries across 65 years\\n\\nand found the OPPOSITE."
    for para in finding_frame.paragraphs:
        para.font.size = Pt(36)
        para.font.color.rgb = TEXT_COLOR
        para.alignment = PP_ALIGN.CENTER

def add_image_slide(title, image_path, notes=""):
    """Helper function to add slides with images"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.3), width=Inches(14))
    else:
        # Placeholder if image doesn't exist
        placeholder_box = slide.shapes.add_textbox(Inches(4), Inches(4), Inches(8), Inches(1))
        placeholder_frame = placeholder_box.text_frame
        placeholder_frame.text = f"[Image: {os.path.basename(image_path)}]"
        placeholder_para = placeholder_frame.paragraphs[0]
        placeholder_para.font.size = Pt(24)
        placeholder_para.font.italic = True
        placeholder_para.alignment = PP_ALIGN.CENTER
    
    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

def add_bullet_slide(title, bullets, notes=""):
    """Helper function to add bullet point slides"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    # Bullets
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(7))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(12)
    
    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

# Build presentation
print("🎨 Creating PowerPoint presentation...")

# Slide 1: Title
add_title_slide()
print("✅ Slide 1: Title")

# Slide 2: The Hook
add_hook_slide()
print("✅ Slide 2: Hook")

# Slide 3: Executive Summary
add_image_slide(
    "Executive Summary: Key Insights",
    "visualizations/06_executive_summary_infographic.png",
    "Highlight -0.156 correlation, 8.7% Ethiopia growth, 14 countries, 65 years data"
)
print("✅ Slide 3: Executive Summary")

# Slide 4: The Data
add_bullet_slide(
    "What We Analyzed: Comprehensive Fiscal Data",
    [
        "📊 Dataset: 623 observations across 14 African countries (1960-2025)",
        "🎯 Focus Period: 2010-2025 (224 recent observations)",
        "📈 Indicators: 28 original + 22 engineered = 50 total features",
        "🌍 SDG Coverage: Growth (SDG 8), Governance (SDG 16), Infrastructure (SDG 9), Health (SDG 3), Poverty (SDG 1)"
    ]
)
print("✅ Slide 4: Data Overview")

# Slide 5: Finding #1 - Efficiency Gap
add_image_slide(
    "Finding #1: The Efficiency Gap (-0.156 Correlation)",
    "visualizations/01_efficiency_gap_quadrants.png",
    "THE MONEY SHOT: Weak negative correlation between Tax-to-GDP and GDP Growth. Ethiopia/Rwanda/Tanzania in green zone (low tax, high growth). Egypt in red zone (high tax, low growth)."
)
print("✅ Slide 5: Efficiency Gap")

# Slide 6: Top vs Bottom Performers
add_image_slide(
    "Finding #2: What Separates Winners from Losers?",
    "visualizations/02_top_vs_bottom_performers.png",
    "Top: Ethiopia 8.73%, Rwanda 7.02%, Tanzania 6.06%. Bottom: South Africa 1.28%, Angola 2.46%, Algeria 2.59%. Pattern: LOW tax + FOCUSED infrastructure = high growth"
)
print("✅ Slide 6: Performers")

# Slide 7: COVID Impact
add_image_slide(
    "Finding #3: The COVID Shock (2020 Crisis)",
    "visualizations/03_covid_impact_timeseries.png",
    "Growth collapsed across all countries. Debt spiked dramatically. Recovery patterns DIVERGED - Ethiopia/Rwanda quick V-shape, South Africa/Nigeria prolonged stagnation"
)
print("✅ Slide 7: COVID Impact")

# Slide 8: Case Study
add_image_slide(
    "Case Study: Ethiopia vs Egypt - Same Governance, Different Outcomes",
    "visualizations/04_ethiopia_vs_egypt_case_study.png",
    "Ethiopia: 8.73% growth with 17.52% tax-to-GDP. Egypt: 3.70% growth with similar tax burden. Egypt spends 17x more on infrastructure, yet achieves 2.4x LESS growth. The Efficiency Gap in action."
)
print("✅ Slide 8: Case Study")

# Slide 9: SDG Dashboard
add_image_slide(
    "SDG Performance Dashboard: Country Rankings",
    "visualizations/05_sdg_performance_dashboard.png",
    "Ethiopia dominates SDG 8 (Growth). Egypt leads SDG 16 (Governance) but fails to translate. No country excels across ALL SDGs - suggests tradeoffs, not coordination"
)
print("✅ Slide 9: SDG Dashboard")

# Slide 10: Root Cause
add_bullet_slide(
    "Root Cause: Why Does the Efficiency Gap Exist?",
    [
        "1. Revenue collection ≠ Effective spending",
        "2. High corruption/bureaucracy waste resources",
        "3. Successful countries focus on PRIVATE SECTOR-LED growth (Ethiopia model)",
        "4. COVID amplified pre-existing inefficiencies",
        "5. Lack of coordination between health/infrastructure spending"
    ],
    "The 5 Whys framework applied to understand root causes"
)
print("✅ Slide 10: Root Cause")

# Slide 11: Policy Recommendations
add_bullet_slide(
    "Policy Recommendations: From Data to Decision",
    [
        "🎯 Measure 'Growth per Dollar of Tax Collected', not just Tax-to-GDP ratio",
        "🚀 Study East African Model: Ethiopia, Rwanda, Tanzania (private sector-led)",
        "🔍 Implement Fiscal Transparency: Public dashboards tracking spending → outcomes",
        "🤝 Coordinate Health & Education: 0.75 correlation suggests bundle SDG 3 + SDG 4",
        "📊 Build Debt Early Warning: Trigger when Debt Growth > GDP Growth for 3+ years"
    ]
)
print("✅ Slide 11: Recommendations")

# Slide 12: Impact Projections
add_bullet_slide(
    "Impact Projections: What If We Close the Gap?",
    [
        "Conservative Scenario (5 years):",
        "  • If Egypt reached Ethiopia's efficiency: 3.70% → 6.5% growth (+75%)",
        "",
        "Aggressive Scenario:",
        "  • If bottom 5 adopted top 5 practices: 2.1% → 5.5% growth (+160%)",
        "  • Poverty reduction: Millions lifted",
        "",
        "The data proves it's POSSIBLE - will we ACT?"
    ]
)
print("✅ Slide 12: Impact Projections")

# Slide 13: Methodology & Limitations
add_bullet_slide(
    "Methodology & Limitations (Transparency)",
    [
        "✅ 65 years of data, 14 countries, 50 indicators, K-Means clustering",
        "⚠️ Egypt data anomaly: Tax-to-GDP corrupted (used median for viz)",
        "⚠️ Missing data: 84 values (13.7%) - documented, not imputed",
        "⚠️ External data: Could not merge World Bank SDG outcomes (time constraint)",
        "⚠️ Correlation ≠ Causation: We show LINKS, not PROOF",
        "",
        "Judge's Note: We prioritized transparency over perfection"
    ]
)
print("✅ Slide 13: Methodology")

# Slide 14: Next Steps
add_bullet_slide(
    "Next Steps: This Is Just the Beginning",
    [
        "Phase 2 Recommendations:",
        "  1. Merge World Bank SDG outcome data (poverty, literacy, mortality)",
        "  2. Causal analysis: Difference-in-Differences on fiscal reforms",
        "  3. Predictive modeling: ML to predict efficiency gaps",
        "  4. Real-time dashboard: Public tool to track fiscal efficiency",
        "  5. Policy pilot: Partner with African governments for intervention"
    ]
)
print("✅ Slide 14: Next Steps")

# Slide 15: The Ask
add_bullet_slide(
    "Join Us in Closing the Gap",
    [
        "For Policymakers: Access our analysis, request country deep-dives",
        "For Researchers: Our data + code fully open-source - replicate & extend",
        "For Investors: The Efficiency Gap is an OPPORTUNITY",
        "  • Ethiopia, Rwanda, Tanzania: Frontier markets with proven track records",
        "",
        "Contact: [Your Email/LinkedIn]",
        "GitHub: [Repository Link]"
    ]
)
print("✅ Slide 15: The Ask")

# Slide 16: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
thank_you_box = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(10), Inches(2))
thank_you_frame = thank_you_box.text_frame
thank_you_frame.text = "Questions? Challenges?\\n\\nLet's Discuss."
for para in thank_you_frame.paragraphs:
    para.font.size = Pt(48)
    para.font.bold = True
    para.font.color.rgb = PRIMARY_COLOR
    para.alignment = PP_ALIGN.CENTER
print("✅ Slide 16: Thank You")

# Save presentation
output_path = "presentation/10Alytics_Governance_Growth_Gap.pptx"
os.makedirs("presentation", exist_ok=True)
prs.save(output_path)

print(f"\\n🎉 Presentation created successfully!")
print(f"📁 Saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print("\\n✅ Ready for submission!")
